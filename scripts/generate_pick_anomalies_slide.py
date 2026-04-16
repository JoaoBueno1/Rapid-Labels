#!/usr/bin/env python3
"""
Pick Anomalies — Clean One-Pager for existing template
======================================================
Designed to integrate into an existing corporate PowerPoint.
Light background, clean layout, focused on:
  - 3 phases accuracy comparison (Before Stocktake / After Stocktake / After Scanners)
  - System key benefits
  - Location correction advantage
  - Target
"""

import os
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ─── Color Palette (light theme matching corporate template) ───
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
NEAR_WHITE   = RGBColor(0xF5, 0xF5, 0xF5)
LIGHT_BG     = RGBColor(0xED, 0xED, 0xED)  # Slide bg
TEXT_DARK    = RGBColor(0x1A, 0x1A, 0x2E)   # Near-black for titles
TEXT_BODY    = RGBColor(0x33, 0x33, 0x40)   # Dark gray for body
TEXT_SUBTLE  = RGBColor(0x6B, 0x6B, 0x7B)   # Subtle gray
ACCENT_BLUE  = RGBColor(0x00, 0x9E, 0xDB)   # Corporate bright blue (matching template stripe)
DARK_ACCENT  = RGBColor(0x2D, 0x2D, 0x3D)   # Dark stripe color
PHASE_RED    = RGBColor(0xDC, 0x35, 0x45)
PHASE_AMBER  = RGBColor(0xF0, 0x9A, 0x00)
PHASE_GREEN  = RGBColor(0x0F, 0x9D, 0x58)
CARD_BG      = RGBColor(0xFA, 0xFA, 0xFA)
CARD_BORDER  = RGBColor(0xDE, 0xDE, 0xDE)
TARGET_BG    = RGBColor(0xE8, 0xF5, 0xE9)
TARGET_BDR   = RGBColor(0x81, 0xC7, 0x84)
CHECK_GREEN  = RGBColor(0x2E, 0x7D, 0x32)


def add_shape(slide, shape_type, left, top, width, height, fill=None, border=None, border_width=Pt(1)):
    s = slide.shapes.add_shape(shape_type, left, top, width, height)
    if fill:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if border:
        s.line.color.rgb = border
        s.line.width = border_width
    else:
        s.line.fill.background()
    return s


def add_text(slide, left, top, width, height, text, size=12, color=TEXT_DARK,
             bold=False, align=PP_ALIGN.LEFT, font="Segoe UI", italic=False):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font
    p.font.italic = italic
    p.alignment = align
    return tb


def add_multiline(slide, left, top, width, height, lines, default_size=10,
                  default_color=TEXT_BODY, font="Segoe UI"):
    """Add textbox with multiple formatted lines.
    Each line: (text, size, color, bold, italic) or just string."""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True

    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        if isinstance(line, str):
            p.text = line
            p.font.size = Pt(default_size)
            p.font.color.rgb = default_color
            p.font.name = font
        else:
            text, size, color, bold, italic = (line + (None, None, None, None, None))[:5]
            p.text = text or ""
            p.font.size = Pt(size or default_size)
            p.font.color.rgb = color or default_color
            p.font.bold = bold or False
            p.font.italic = italic or False
            p.font.name = font
        p.space_after = Pt(2)

    return tb


def build_slide():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # ── NO background shape — let it be default white/light so it blends
    # with the existing template when user copies the slide ──

    SW = prs.slide_width
    SH = prs.slide_height

    # ════════════════════════════════════════════
    # TITLE AREA (top)
    # ════════════════════════════════════════════

    # Blue accent line under title
    add_shape(slide, MSO_SHAPE.RECTANGLE,
              Inches(0.8), Inches(0.95), Inches(1.8), Inches(0.045),
              fill=ACCENT_BLUE)

    add_text(slide, Inches(0.8), Inches(0.3), Inches(10), Inches(0.5),
             "Pick Anomalies System", size=26, color=TEXT_DARK, bold=True)

    add_text(slide, Inches(0.8), Inches(1.05), Inches(10), Inches(0.35),
             "Real-time pick location detection — identifying and correcting warehouse stock discrepancies",
             size=11, color=TEXT_SUBTLE, italic=True)

    # ════════════════════════════════════════════
    # THE 3 PHASES (main visual — big numbers)
    # ════════════════════════════════════════════

    phases_y = Inches(1.65)

    # Section label
    add_text(slide, Inches(0.8), phases_y - Inches(0.02), Inches(5), Inches(0.25),
             "PICK ACCURACY EVOLUTION", size=9, color=ACCENT_BLUE, bold=True)

    card_w = Inches(3.2)
    card_h = Inches(2.35)
    card_gap = Inches(0.4)
    cards_start_x = Inches(0.8)

    phases = [
        {
            "title": "Before Stocktake",
            "period": "Pre-February 2026",
            "pct": "~78%",
            "color": PHASE_RED,
            "bar_pct": 0.78,
            "bullets": [
                "No visibility into pick errors",
                "Discrepancies found only at stocktake",
                "Stock variance unknown daily",
            ]
        },
        {
            "title": "After Stocktake",
            "period": "February – March 2026",
            "pct": "~88%",
            "color": PHASE_AMBER,
            "bar_pct": 0.88,
            "bullets": [
                "Bin locations corrected via full count",
                "Improved but errors still invisible",
                "No way to detect daily pick issues",
            ]
        },
        {
            "title": "Scanners + Pick Anomalies",
            "period": "March 26 – Present",
            "pct": "95.4%",
            "color": PHASE_GREEN,
            "bar_pct": 0.954,
            "bullets": [
                "Anomalies detected within hours",
                "1,076 orders analyzed automatically",
                "105 anomalies found, 62 corrected",
            ]
        },
    ]

    for i, ph in enumerate(phases):
        cx = cards_start_x + i * (card_w + card_gap)

        # Card background
        card = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE,
                         cx, phases_y + Inches(0.28), card_w, card_h,
                         fill=CARD_BG, border=CARD_BORDER)
        card.adjustments[0] = 0.04

        # Color bar at top of card
        add_shape(slide, MSO_SHAPE.RECTANGLE,
                  cx + Inches(0.01), phases_y + Inches(0.29),
                  card_w - Inches(0.02), Inches(0.055),
                  fill=ph["color"])

        # Title
        add_text(slide, cx + Inches(0.2), phases_y + Inches(0.42),
                 card_w - Inches(0.4), Inches(0.25),
                 ph["title"], size=12, color=TEXT_DARK, bold=True)

        # Period
        add_text(slide, cx + Inches(0.2), phases_y + Inches(0.68),
                 card_w - Inches(0.4), Inches(0.2),
                 ph["period"], size=8, color=TEXT_SUBTLE)

        # Big percentage
        add_text(slide, cx + Inches(0.2), phases_y + Inches(0.90),
                 card_w - Inches(0.4), Inches(0.55),
                 ph["pct"], size=38, color=ph["color"], bold=True,
                 align=PP_ALIGN.CENTER)

        # Progress bar background
        bar_y = phases_y + Inches(1.50)
        bar_full_w = card_w - Inches(0.4)
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE,
                  cx + Inches(0.2), bar_y,
                  bar_full_w, Inches(0.09),
                  fill=RGBColor(0xE0, 0xE0, 0xE0))
        # Progress bar fill
        fill_w = int(bar_full_w * ph["bar_pct"])
        if fill_w > 0:
            bar_fill = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE,
                                 cx + Inches(0.2), bar_y,
                                 fill_w, Inches(0.09),
                                 fill=ph["color"])

        # Bullet points
        for j, bullet in enumerate(ph["bullets"]):
            by = phases_y + Inches(1.70) + j * Inches(0.20)
            add_text(slide, cx + Inches(0.2), by,
                     Inches(0.15), Inches(0.18),
                     "•", size=8, color=ph["color"], bold=True)
            add_text(slide, cx + Inches(0.38), by,
                     card_w - Inches(0.58), Inches(0.18),
                     bullet, size=8, color=TEXT_BODY)

    # Arrow between cards
    for i in range(2):
        ax = cards_start_x + (i + 1) * (card_w + card_gap) - card_gap + Inches(0.03)
        arrow = add_shape(slide, MSO_SHAPE.RIGHT_ARROW,
                          ax, phases_y + Inches(1.3),
                          Inches(0.34), Inches(0.22),
                          fill=RGBColor(0xBB, 0xBB, 0xBB))

    # Improvement callout on far right
    imp_x = cards_start_x + 3 * (card_w + card_gap) - Inches(0.15)
    imp_y = phases_y + Inches(0.70)
    imp_box = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE,
                        imp_x, imp_y, Inches(1.65), Inches(1.5),
                        fill=TARGET_BG, border=TARGET_BDR)
    imp_box.adjustments[0] = 0.08

    add_text(slide, imp_x, imp_y + Inches(0.12),
             Inches(1.65), Inches(0.55),
             "+17.4%", size=30, color=CHECK_GREEN, bold=True,
             align=PP_ALIGN.CENTER)
    add_text(slide, imp_x, imp_y + Inches(0.70),
             Inches(1.65), Inches(0.40),
             "Overall\nImprovement", size=10, color=CHECK_GREEN,
             align=PP_ALIGN.CENTER)
    add_text(slide, imp_x, imp_y + Inches(1.1),
             Inches(1.65), Inches(0.30),
             "78% → 95.4%", size=9, color=TEXT_SUBTLE,
             align=PP_ALIGN.CENTER)

    # ════════════════════════════════════════════
    # BOTTOM SECTION: System Benefits + Target
    # ════════════════════════════════════════════

    bottom_y = Inches(4.45)

    # Subtle divider
    add_shape(slide, MSO_SHAPE.RECTANGLE,
              Inches(0.8), bottom_y - Inches(0.06),
              Inches(11.7), Inches(0.015),
              fill=RGBColor(0xD0, 0xD0, 0xD0))

    # ── Left: System Key Benefits ──
    add_text(slide, Inches(0.8), bottom_y + Inches(0.05),
             Inches(4), Inches(0.25),
             "WHAT THE SYSTEM DOES", size=9, color=ACCENT_BLUE, bold=True)

    benefits = [
        ("Real-Time Anomaly Detection",
         "Every shipped order is analyzed automatically. When a picker takes stock from the wrong bin, the system flags it within hours — not weeks."),
        ("Automatic Stock Corrections",
         "Once reviewed, the system creates Stock Transfers in Cin7 to move units back to the correct location. 62 corrections executed so far."),
        ("Bin Location Maintenance",
         "The system continuously identifies products stored in wrong locations, helping the team relocate items and keep bin assignments accurate across the warehouse."),
        ("Stock Variance Prevention",
         "By catching discrepancies early — before they compound — we prevent significant stock variances that previously were only discovered during full stocktakes."),
    ]

    for i, (title, desc) in enumerate(benefits):
        by = bottom_y + Inches(0.38) + i * Inches(0.62)

        # Check icon circle
        circle = add_shape(slide, MSO_SHAPE.OVAL,
                           Inches(0.85), by + Inches(0.02),
                           Inches(0.22), Inches(0.22),
                           fill=PHASE_GREEN)
        add_text(slide, Inches(0.85), by + Inches(0.01),
                 Inches(0.22), Inches(0.22),
                 "✓", size=10, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

        # Benefit title
        add_text(slide, Inches(1.18), by,
                 Inches(5), Inches(0.22),
                 title, size=10, color=TEXT_DARK, bold=True)

        # Description
        add_text(slide, Inches(1.18), by + Inches(0.22),
                 Inches(5.3), Inches(0.35),
                 desc, size=8, color=TEXT_SUBTLE)

    # ── Right: Live Numbers + Target ──
    right_x = Inches(7.2)

    # Live stats card
    add_text(slide, right_x, bottom_y + Inches(0.05),
             Inches(3), Inches(0.25),
             "LIVE DATA", size=9, color=ACCENT_BLUE, bold=True)

    stats_card = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE,
                           right_x, bottom_y + Inches(0.35),
                           Inches(2.7), Inches(1.35),
                           fill=CARD_BG, border=CARD_BORDER)
    stats_card.adjustments[0] = 0.05

    stats = [
        ("1,076", "Orders Analyzed"),
        ("2,286", "Total Picks Tracked"),
        ("105", "Anomalies Detected"),
        ("62", "Corrections Executed"),
        ("95.4%", "Current Pick Accuracy"),
    ]

    for i, (val, label) in enumerate(stats):
        sy = bottom_y + Inches(0.42) + i * Inches(0.23)
        vc = PHASE_GREEN if "95" in val else ACCENT_BLUE
        add_text(slide, right_x + Inches(0.15), sy,
                 Inches(0.8), Inches(0.22),
                 val, size=11, color=vc, bold=True, align=PP_ALIGN.RIGHT)
        add_text(slide, right_x + Inches(1.05), sy,
                 Inches(1.5), Inches(0.22),
                 label, size=9, color=TEXT_BODY)

    # Target card
    target_y = bottom_y + Inches(1.85)
    add_text(slide, right_x, target_y - Inches(0.02),
             Inches(3), Inches(0.25),
             "TARGET", size=9, color=ACCENT_BLUE, bold=True)

    tgt_card = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE,
                         right_x, target_y + Inches(0.22),
                         Inches(2.7), Inches(0.68),
                         fill=TARGET_BG, border=TARGET_BDR)
    tgt_card.adjustments[0] = 0.06

    add_text(slide, right_x + Inches(0.1), target_y + Inches(0.28),
             Inches(1.3), Inches(0.50),
             "97–98%", size=26, color=CHECK_GREEN, bold=True)

    add_text(slide, right_x + Inches(1.45), target_y + Inches(0.30),
             Inches(1.2), Inches(0.55),
             "Pick Accuracy\nGoal by Q3 2026\nwith process improvements",
             size=7.5, color=CHECK_GREEN)

    # ── Far right: Improvement roadmap ──
    road_x = Inches(10.2)

    add_text(slide, road_x, bottom_y + Inches(0.05),
             Inches(3), Inches(0.25),
             "NEXT STEPS", size=9, color=ACCENT_BLUE, bold=True)

    next_items = [
        "Real-time stock updates via Cin7 Webhooks",
        "Repeat-offender SKU relocation program",
        "Warehouse team process training",
        "Bin audit scheduling based on anomaly data",
        "Reduce stock variance before it impacts orders",
    ]

    for i, item in enumerate(next_items):
        ny = bottom_y + Inches(0.40) + i * Inches(0.35)
        # Number circle
        nc = add_shape(slide, MSO_SHAPE.OVAL,
                       road_x, ny + Inches(0.01),
                       Inches(0.20), Inches(0.20),
                       fill=ACCENT_BLUE)
        add_text(slide, road_x, ny,
                 Inches(0.20), Inches(0.22),
                 str(i + 1), size=8, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, road_x + Inches(0.28), ny,
                 Inches(2.7), Inches(0.30),
                 item, size=8.5, color=TEXT_BODY)

    return prs


def main():
    output_dir = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))), "reports")
    os.makedirs(output_dir, exist_ok=True)

    date_str = datetime.now().strftime("%d%b%Y")
    output_path = os.path.join(output_dir, f"Pick_Anomalies_Slide_{date_str}.pptx")

    print("🎨 Building Pick Anomalies slide (template-compatible)...")
    prs = build_slide()
    prs.save(output_path)
    print(f"✅ Saved to: {output_path}")
    return output_path


if __name__ == "__main__":
    path = main()
    print(path)
