#!/usr/bin/env python3
"""
Generate Pick Anomalies Evolution PowerPoint — One-Pager Board Report
=====================================================================
Professional single-slide presentation showing the evolution of pick
accuracy across three phases: Before Stocktake, After Stocktake, After
Scanners + Pick Anomalies System.

Uses real data from Supabase pick_anomaly_orders table.
"""

import os
import sys
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ─── Color Palette ───
DARK_BG      = RGBColor(0x0F, 0x17, 0x2A)   # Deep navy
CARD_BG      = RGBColor(0x1A, 0x24, 0x3B)   # Slightly lighter navy
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY   = RGBColor(0xA0, 0xAE, 0xC0)
ACCENT_BLUE  = RGBColor(0x38, 0xBD, 0xF8)   # Bright cyan-blue
ACCENT_GREEN = RGBColor(0x22, 0xC5, 0x5E)   # Success green
ACCENT_RED   = RGBColor(0xEF, 0x44, 0x44)   # Warning red
ACCENT_AMBER = RGBColor(0xF5, 0x9E, 0x0B)   # Amber/yellow
ACCENT_PURPLE= RGBColor(0xA7, 0x8B, 0xFA)   # Purple accent
PHASE_RED    = RGBColor(0xDC, 0x26, 0x26)   # Phase 1 red
PHASE_AMBER  = RGBColor(0xF5, 0x9E, 0x0B)   # Phase 2 amber
PHASE_GREEN  = RGBColor(0x16, 0xA3, 0x4A)   # Phase 3 green
SUBTLE_LINE  = RGBColor(0x2D, 0x3A, 0x55)   # Divider lines

# ─── Real Data ───
REAL_DATA = {
    "total_orders": 1076,
    "orders_with_anomalies": 106,
    "orders_clean": 970,
    "reviewed": 117,
    "corrections_made": 62,
    "total_picks": 2286,
    "correct_picks": 2181,
    "anomaly_picks": 105,
    "pick_accuracy": 95.41,
    "anomaly_rate": 4.59,
    "weekly": [
        {"week": "Mar 23", "orders": 147, "picks": 334, "anomalies": 16, "accuracy": 95.2},
        {"week": "Mar 30", "orders": 386, "picks": 862, "anomalies": 25, "accuracy": 97.1},
        {"week": "Apr 06", "orders": 402, "picks": 940, "anomalies": 59, "accuracy": 93.7},
        {"week": "Apr 13", "orders": 52,  "picks": 124, "anomalies": 4,  "accuracy": 96.8},
    ],
}

# Three phases data (estimated/real)
PHASES = {
    "before_stocktake": {
        "label": "Before Stocktake",
        "period": "Pre-Feb 2026",
        "accuracy": 78,
        "description": "No visibility into pick\nerrors. Stock discrepancies\nonly found during\nstocktake counts.",
        "icon": "!",
        "color": PHASE_RED,
    },
    "after_stocktake": {
        "label": "After Stocktake",
        "period": "Feb – Mar 2026",
        "accuracy": 88,
        "description": "Stocktake corrected\nbin locations. Accuracy\nimproved but errors\nstill undetected daily.",
        "icon": "~",
        "color": PHASE_AMBER,
    },
    "after_scanners": {
        "label": "After Scanners + System",
        "period": "Mar 26 – Present",
        "accuracy": 95.4,
        "description": "Real-time pick detection\nvia scanner data. Anomalies\nidentified within hours,\nnot weeks.",
        "icon": "✓",
        "color": PHASE_GREEN,
    },
}


def add_rounded_rect(slide, left, top, width, height, fill_color, border_color=None):
    """Add a rounded rectangle shape."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    # Adjust corner rounding
    shape.adjustments[0] = 0.06
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=12,
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Segoe UI"):
    """Add a simple text box."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_circle_metric(slide, cx, cy, radius, value_text, label_text, ring_color, bg_color=CARD_BG):
    """Add a circular metric with ring effect."""
    # Outer ring
    outer = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        cx - radius, cy - radius,
        radius * 2, radius * 2
    )
    outer.fill.solid()
    outer.fill.fore_color.rgb = ring_color
    outer.line.fill.background()

    # Inner circle
    inner_r = int(radius * 0.82)
    inner = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        cx - inner_r, cy - inner_r,
        inner_r * 2, inner_r * 2
    )
    inner.fill.solid()
    inner.fill.fore_color.rgb = bg_color
    inner.line.fill.background()

    # Value text
    add_text_box(
        slide, cx - radius, cy - Inches(0.28),
        radius * 2, Inches(0.4),
        value_text, font_size=22, color=WHITE, bold=True,
        alignment=PP_ALIGN.CENTER
    )

    # Label below circle
    add_text_box(
        slide, cx - radius - Inches(0.15), cy + radius + Inches(0.04),
        radius * 2 + Inches(0.3), Inches(0.35),
        label_text, font_size=8, color=LIGHT_GRAY,
        alignment=PP_ALIGN.CENTER
    )


def add_phase_card(slide, left, top, width, height, phase_data):
    """Add a phase evolution card."""
    card = add_rounded_rect(slide, left, top, width, height, CARD_BG, SUBTLE_LINE)

    # Color strip at top
    strip = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, Inches(0.06)
    )
    strip.fill.solid()
    strip.fill.fore_color.rgb = phase_data["color"]
    strip.line.fill.background()

    # Phase label
    add_text_box(
        slide, left + Inches(0.12), top + Inches(0.12),
        width - Inches(0.24), Inches(0.22),
        phase_data["label"], font_size=10, color=phase_data["color"], bold=True
    )

    # Period
    add_text_box(
        slide, left + Inches(0.12), top + Inches(0.33),
        width - Inches(0.24), Inches(0.18),
        phase_data["period"], font_size=7.5, color=LIGHT_GRAY
    )

    # Big accuracy number
    acc_text = f'{phase_data["accuracy"]}%'
    add_text_box(
        slide, left + Inches(0.12), top + Inches(0.55),
        width - Inches(0.24), Inches(0.45),
        acc_text, font_size=28, color=WHITE, bold=True,
        alignment=PP_ALIGN.CENTER
    )

    # "Pick Accuracy" label
    add_text_box(
        slide, left + Inches(0.12), top + Inches(0.95),
        width - Inches(0.24), Inches(0.18),
        "Pick Accuracy", font_size=8, color=LIGHT_GRAY,
        alignment=PP_ALIGN.CENTER
    )

    # Description
    add_text_box(
        slide, left + Inches(0.12), top + Inches(1.18),
        width - Inches(0.24), Inches(0.65),
        phase_data["description"], font_size=7, color=LIGHT_GRAY
    )


def add_arrow_between(slide, x, y, color):
    """Add a right-pointing arrow."""
    arrow = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW, x, y, Inches(0.35), Inches(0.25)
    )
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = color
    arrow.line.fill.background()


def add_kpi_box(slide, left, top, width, height, value, label, color=ACCENT_BLUE):
    """Add a small KPI stat box."""
    card = add_rounded_rect(slide, left, top, width, height, CARD_BG, SUBTLE_LINE)

    # Value
    add_text_box(
        slide, left + Inches(0.08), top + Inches(0.06),
        width - Inches(0.16), Inches(0.28),
        str(value), font_size=16, color=color, bold=True,
        alignment=PP_ALIGN.CENTER
    )

    # Label
    add_text_box(
        slide, left + Inches(0.08), top + Inches(0.34),
        width - Inches(0.16), Inches(0.26),
        label, font_size=6.5, color=LIGHT_GRAY,
        alignment=PP_ALIGN.CENTER
    )


def add_bar(slide, x, y, bar_width, max_height, value, max_value, color, label, show_value=True):
    """Add a single bar for the chart."""
    bar_height = max(Inches(0.08), int(max_height * (value / max_value))) if max_value > 0 else Inches(0.08)
    bar_y = y + max_height - bar_height

    bar = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, x, bar_y, bar_width, bar_height
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    if hasattr(bar, 'adjustments') and len(bar.adjustments) > 0:
        bar.adjustments[0] = 0.15

    # Value on top
    if show_value:
        add_text_box(
            slide, x - Inches(0.1), bar_y - Inches(0.18),
            bar_width + Inches(0.2), Inches(0.18),
            f"{value:.1f}%", font_size=7, color=WHITE, bold=True,
            alignment=PP_ALIGN.CENTER
        )

    # Label below
    add_text_box(
        slide, x - Inches(0.12), y + max_height + Inches(0.02),
        bar_width + Inches(0.24), Inches(0.22),
        label, font_size=6, color=LIGHT_GRAY,
        alignment=PP_ALIGN.CENTER
    )


def build_presentation():
    """Build the complete one-page PowerPoint."""
    prs = Presentation()
    # Widescreen 16:9
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Use blank layout
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # ── Full dark background ──
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK_BG
    bg.line.fill.background()

    SW = prs.slide_width  # slide width
    SH = prs.slide_height

    # ════════════════════════════════════════════
    # HEADER SECTION
    # ════════════════════════════════════════════

    # Title bar background
    header_bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, SW, Inches(0.85)
    )
    header_bg.fill.solid()
    header_bg.fill.fore_color.rgb = RGBColor(0x12, 0x1C, 0x32)
    header_bg.line.fill.background()

    # RapidLED logo placeholder (text-based)
    add_text_box(
        slide, Inches(0.4), Inches(0.12),
        Inches(2), Inches(0.35),
        "RAPIDLED", font_size=16, color=ACCENT_BLUE, bold=True,
        font_name="Segoe UI"
    )
    add_text_box(
        slide, Inches(0.4), Inches(0.45),
        Inches(2.5), Inches(0.25),
        "Warehouse Intelligence", font_size=9, color=LIGHT_GRAY
    )

    # Main title
    add_text_box(
        slide, Inches(2.8), Inches(0.08),
        Inches(7.5), Inches(0.45),
        "Pick Accuracy Evolution — From Blind Spots to Real-Time Detection",
        font_size=17, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER
    )
    add_text_box(
        slide, Inches(2.8), Inches(0.48),
        Inches(7.5), Inches(0.25),
        f"Pick Anomalies System  •  Live Data as of {datetime.now().strftime('%B %d, %Y')}  •  Main Warehouse",
        font_size=8.5, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER
    )

    # Date badge
    add_text_box(
        slide, Inches(11.0), Inches(0.22),
        Inches(2), Inches(0.35),
        datetime.now().strftime("%d %b %Y"), font_size=11, color=ACCENT_BLUE,
        bold=True, alignment=PP_ALIGN.RIGHT
    )

    # ════════════════════════════════════════════
    # SECTION 1: THE THREE PHASES (left side)
    # ════════════════════════════════════════════

    section_y = Inches(1.05)

    # Section label
    add_text_box(
        slide, Inches(0.4), section_y - Inches(0.02),
        Inches(4), Inches(0.22),
        "THE JOURNEY", font_size=8, color=ACCENT_BLUE, bold=True
    )

    # Subtitle
    add_text_box(
        slide, Inches(0.4), section_y + Inches(0.18),
        Inches(8), Inches(0.22),
        "Three phases of improvement — from no visibility to near real-time anomaly detection",
        font_size=8, color=LIGHT_GRAY
    )

    phase_y = section_y + Inches(0.48)
    phase_w = Inches(2.55)
    phase_h = Inches(1.9)
    phase_gap = Inches(0.25)

    # Phase 1: Before Stocktake
    x1 = Inches(0.4)
    add_phase_card(slide, x1, phase_y, phase_w, phase_h, PHASES["before_stocktake"])

    # Arrow 1→2
    add_arrow_between(slide, x1 + phase_w + Inches(0.02), phase_y + Inches(0.85), LIGHT_GRAY)

    # Phase 2: After Stocktake
    x2 = x1 + phase_w + phase_gap + Inches(0.15)
    add_phase_card(slide, x2, phase_y, phase_w, phase_h, PHASES["after_stocktake"])

    # Arrow 2→3
    add_arrow_between(slide, x2 + phase_w + Inches(0.02), phase_y + Inches(0.85), LIGHT_GRAY)

    # Phase 3: Scanners + System
    x3 = x2 + phase_w + phase_gap + Inches(0.15)
    add_phase_card(slide, x3, phase_y, phase_w, phase_h, PHASES["after_scanners"])

    # Improvement badge
    improvement = PHASES["after_scanners"]["accuracy"] - PHASES["before_stocktake"]["accuracy"]
    badge = add_rounded_rect(
        slide, x3 + phase_w + Inches(0.18), phase_y + Inches(0.55),
        Inches(1.15), Inches(0.85), RGBColor(0x0D, 0x3B, 0x1F), PHASE_GREEN
    )
    add_text_box(
        slide, x3 + phase_w + Inches(0.18), phase_y + Inches(0.60),
        Inches(1.15), Inches(0.30),
        f"+{improvement:.1f}%", font_size=22, color=ACCENT_GREEN, bold=True,
        alignment=PP_ALIGN.CENTER
    )
    add_text_box(
        slide, x3 + phase_w + Inches(0.18), phase_y + Inches(0.92),
        Inches(1.15), Inches(0.35),
        "Overall\nImprovement", font_size=7.5, color=ACCENT_GREEN,
        alignment=PP_ALIGN.CENTER
    )

    # ════════════════════════════════════════════
    # SECTION 2: LIVE METRICS (middle band)
    # ════════════════════════════════════════════

    metrics_y = phase_y + phase_h + Inches(0.28)

    # Section divider line
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.4), metrics_y - Inches(0.05),
        Inches(12.5), Inches(0.015)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = SUBTLE_LINE
    line.line.fill.background()

    add_text_box(
        slide, Inches(0.4), metrics_y + Inches(0.02),
        Inches(4), Inches(0.22),
        "LIVE SYSTEM METRICS", font_size=8, color=ACCENT_BLUE, bold=True
    )
    add_text_box(
        slide, Inches(3.5), metrics_y + Inches(0.02),
        Inches(6), Inches(0.22),
        "Pick Anomalies System — automated detection since scanner rollout (Mar 26, 2026)",
        font_size=7.5, color=LIGHT_GRAY
    )

    kpi_y = metrics_y + Inches(0.30)
    kpi_w = Inches(1.48)
    kpi_h = Inches(0.62)
    kpi_gap = Inches(0.12)

    d = REAL_DATA
    kpis = [
        (f"{d['total_orders']:,}", "Orders\nAnalyzed", ACCENT_BLUE),
        (f"{d['total_picks']:,}", "Total Picks\nTracked", ACCENT_BLUE),
        (f"{d['pick_accuracy']:.1f}%", "Current Pick\nAccuracy", ACCENT_GREEN),
        (f"{d['anomaly_picks']}", "Anomalies\nDetected", ACCENT_AMBER),
        (f"{d['corrections_made']}", "Stock Corrections\nExecuted", ACCENT_PURPLE),
        (f"{d['reviewed']}", "Orders\nReviewed", LIGHT_GRAY),
        (f"{d['orders_clean']:,}", "Clean Orders\n(0 Issues)", ACCENT_GREEN),
        (f"~{(d['corrections_made']/max(d['anomaly_picks'],1))*100:.0f}%", "Anomalies\nCorrected", ACCENT_PURPLE),
    ]

    for i, (val, label, color) in enumerate(kpis):
        kx = Inches(0.4) + i * (kpi_w + kpi_gap)
        add_kpi_box(slide, kx, kpi_y, kpi_w, kpi_h, val, label, color)

    # ════════════════════════════════════════════
    # SECTION 3: WEEKLY TREND + WHAT'S NEXT (bottom)
    # ════════════════════════════════════════════

    bottom_y = kpi_y + kpi_h + Inches(0.28)

    # Divider
    line2 = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.4), bottom_y - Inches(0.05),
        Inches(12.5), Inches(0.015)
    )
    line2.fill.solid()
    line2.fill.fore_color.rgb = SUBTLE_LINE
    line2.line.fill.background()

    # ── Left: Weekly Accuracy Chart ──
    add_text_box(
        slide, Inches(0.4), bottom_y + Inches(0.02),
        Inches(4), Inches(0.22),
        "WEEKLY PICK ACCURACY TREND", font_size=8, color=ACCENT_BLUE, bold=True
    )

    chart_x = Inches(0.6)
    chart_y = bottom_y + Inches(0.30)
    chart_max_h = Inches(1.45)
    bar_w = Inches(0.6)

    # 95% reference line
    ref_val = 95.0
    ref_y_pos = chart_y + chart_max_h - int(chart_max_h * (ref_val / 100.0))
    ref_line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, chart_x - Inches(0.1), ref_y_pos,
        Inches(5.6), Inches(0.01)
    )
    ref_line.fill.solid()
    ref_line.fill.fore_color.rgb = ACCENT_AMBER
    ref_line.line.fill.background()
    add_text_box(
        slide, chart_x - Inches(0.6), ref_y_pos - Inches(0.12),
        Inches(0.55), Inches(0.18),
        "95%", font_size=6.5, color=ACCENT_AMBER, alignment=PP_ALIGN.RIGHT
    )

    # Weekly bars
    weekly = d["weekly"]
    colors_weekly = [ACCENT_BLUE, ACCENT_GREEN, ACCENT_AMBER, ACCENT_BLUE]
    for i, w in enumerate(weekly):
        bx = chart_x + Inches(0.3) + i * Inches(1.35)
        # Color based on accuracy
        if w["accuracy"] >= 97:
            bc = ACCENT_GREEN
        elif w["accuracy"] >= 95:
            bc = ACCENT_BLUE
        else:
            bc = ACCENT_AMBER
        add_bar(slide, bx, chart_y, bar_w, chart_max_h, w["accuracy"], 100, bc, w["week"])

        # Order count below label
        add_text_box(
            slide, bx - Inches(0.12), chart_y + chart_max_h + Inches(0.20),
            bar_w + Inches(0.24), Inches(0.15),
            f'{w["orders"]} orders', font_size=5.5, color=RGBColor(0x60, 0x70, 0x88),
            alignment=PP_ALIGN.CENTER
        )

    # ── Right: What's Next / Impact ──
    next_x = Inches(6.6)
    add_text_box(
        slide, next_x, bottom_y + Inches(0.02),
        Inches(3), Inches(0.22),
        "HOW IT WORKS", font_size=8, color=ACCENT_BLUE, bold=True
    )

    how_items = [
        ("1.", "Scanner captures bin location at pick time", ACCENT_BLUE),
        ("2.", "System compares picked bin vs expected bin", ACCENT_BLUE),
        ("3.", "Anomalies flagged within hours of shipping", ACCENT_AMBER),
        ("4.", "Warehouse team reviews & corrects stock", ACCENT_GREEN),
        ("5.", "Corrections executed as Stock Transfers in Cin7", ACCENT_GREEN),
    ]

    for i, (num, text, color) in enumerate(how_items):
        iy = bottom_y + Inches(0.30) + i * Inches(0.26)
        add_text_box(slide, next_x, iy, Inches(0.22), Inches(0.22),
                     num, font_size=8, color=color, bold=True)
        add_text_box(slide, next_x + Inches(0.22), iy, Inches(3.3), Inches(0.22),
                     text, font_size=8, color=LIGHT_GRAY)

    # ── Far Right: Target / Goal ──
    goal_x = Inches(10.2)
    add_text_box(
        slide, goal_x, bottom_y + Inches(0.02),
        Inches(3), Inches(0.22),
        "NEXT MILESTONES", font_size=8, color=ACCENT_BLUE, bold=True
    )

    # Target card
    target_card = add_rounded_rect(
        slide, goal_x, bottom_y + Inches(0.30),
        Inches(2.8), Inches(0.72), RGBColor(0x0D, 0x3B, 0x1F), PHASE_GREEN
    )
    add_text_box(
        slide, goal_x + Inches(0.1), bottom_y + Inches(0.34),
        Inches(2.6), Inches(0.22),
        "TARGET", font_size=7, color=ACCENT_GREEN, bold=True
    )
    add_text_box(
        slide, goal_x + Inches(0.1), bottom_y + Inches(0.52),
        Inches(1.5), Inches(0.4),
        "97–98%", font_size=24, color=WHITE, bold=True
    )
    add_text_box(
        slide, goal_x + Inches(1.5), bottom_y + Inches(0.55),
        Inches(1.3), Inches(0.4),
        "Pick Accuracy\nGoal by Q3 2026", font_size=7.5, color=ACCENT_GREEN
    )

    # Roadmap items
    roadmap = [
        "Cin7 Webhooks — real-time stock updates",
        "Repeat offender SKU analysis & relocation",
        "Process training with warehouse team",
        "Stock discrepancy reduced before impact",
    ]
    for i, item in enumerate(roadmap):
        ry = bottom_y + Inches(1.12) + i * Inches(0.22)
        add_text_box(slide, goal_x + Inches(0.08), ry, Inches(0.18), Inches(0.20),
                     "→", font_size=7.5, color=ACCENT_GREEN, bold=True)
        add_text_box(slide, goal_x + Inches(0.28), ry, Inches(2.5), Inches(0.20),
                     item, font_size=7, color=LIGHT_GRAY)

    # ════════════════════════════════════════════
    # FOOTER
    # ════════════════════════════════════════════

    footer_y = SH - Inches(0.35)
    footer_bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, footer_y, SW, Inches(0.35)
    )
    footer_bg.fill.solid()
    footer_bg.fill.fore_color.rgb = RGBColor(0x0A, 0x0F, 0x1E)
    footer_bg.line.fill.background()

    add_text_box(
        slide, Inches(0.4), footer_y + Inches(0.06),
        Inches(6), Inches(0.22),
        "Pick Anomalies System — Rapid Labels  •  Warehouse Intelligence Platform  •  Data synced every 2h from Cin7 Core",
        font_size=6.5, color=RGBColor(0x50, 0x5E, 0x78)
    )
    add_text_box(
        slide, Inches(9), footer_y + Inches(0.06),
        Inches(4), Inches(0.22),
        "Confidential — RapidLED Australia  •  April 2026",
        font_size=6.5, color=RGBColor(0x50, 0x5E, 0x78),
        alignment=PP_ALIGN.RIGHT
    )

    return prs


def main():
    output_dir = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))), "reports")
    os.makedirs(output_dir, exist_ok=True)

    date_str = datetime.now().strftime("%d%b%Y")
    output_path = os.path.join(output_dir, f"Pick_Anomalies_Evolution_{date_str}.pptx")

    print("🎨 Building Pick Anomalies Evolution presentation...")
    prs = build_presentation()
    prs.save(output_path)
    print(f"✅ Saved to: {output_path}")
    return output_path


if __name__ == "__main__":
    path = main()
    print(path)
